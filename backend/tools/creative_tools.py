"""Creative and presentation tools for The Boss."""
import os
from pathlib import Path
from crewai.tools import BaseTool
from pydantic import BaseModel, Field


WORKSPACE = Path(os.getenv("BOSS_WORKSPACE", str(Path.home() / "Desktop" / "EASY-Workspace")))


def _ensure_workspace() -> Path:
    WORKSPACE.mkdir(parents=True, exist_ok=True)
    return WORKSPACE


class PresentationSlide(BaseModel):
    title: str = Field(description="Slide title")
    content: str = Field(description="Slide body text (use newlines for bullet points)")
    notes: str = Field(default="", description="Speaker notes for the slide")


class CreatePresentationInput(BaseModel):
    filename: str = Field(description="Filename for the presentation (e.g., 'campaign-deck.pptx')")
    title: str = Field(description="Presentation title shown on the cover slide")
    subtitle: str = Field(default="", description="Subtitle for the cover slide")
    slides: list[PresentationSlide] = Field(description="List of slides to create")
    theme: str = Field(default="dark", description="Theme: 'dark' or 'light'")


class CreatePresentationTool(BaseTool):
    name: str = "create_presentation"
    description: str = "Create a PowerPoint (.pptx) presentation with slides. Saves to the EASY workspace and opens automatically."
    args_schema: type[BaseModel] = CreatePresentationInput

    def _run(self, filename: str, title: str, subtitle: str = "", slides: list = None, theme: str = "dark") -> str:
        try:
            from pptx import Presentation
            from pptx.util import Inches, Pt, Emu
            from pptx.dml.color import RGBColor
            from pptx.enum.text import PP_ALIGN

            slides = slides or []
            ws = _ensure_workspace()
            if not filename.endswith(".pptx"):
                filename += ".pptx"
            filepath = ws / filename

            prs = Presentation()
            prs.slide_width = Inches(13.33)
            prs.slide_height = Inches(7.5)

            dark = theme == "dark"
            bg_color = RGBColor(0x0F, 0x17, 0x2A) if dark else RGBColor(0xFF, 0xFF, 0xFF)
            title_color = RGBColor(0x34, 0xD3, 0x99) if dark else RGBColor(0x05, 0x96, 0x69)
            text_color = RGBColor(0xE2, 0xE8, 0xF0) if dark else RGBColor(0x1E, 0x29, 0x3B)
            accent_color = RGBColor(0x10, 0xB9, 0x81) if dark else RGBColor(0x06, 0xB6, 0xD4)

            def set_bg(slide):
                fill = slide.background.fill
                fill.solid()
                fill.fore_color.rgb = bg_color

            def add_text(tf, text, size, bold=False, color=None, align=PP_ALIGN.LEFT):
                tf.word_wrap = True
                for para in tf.paragraphs:
                    para.alignment = align
                    for run in para.runs:
                        run.font.size = Pt(size)
                        run.font.bold = bold
                        run.font.color.rgb = color or text_color

            # Cover slide
            cover_layout = prs.slide_layouts[6]  # blank
            slide = prs.slides.add_slide(cover_layout)
            set_bg(slide)

            # Title text box
            txb = slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(11), Inches(1.5))
            tf = txb.text_frame
            tf.text = title
            p = tf.paragraphs[0]
            p.alignment = PP_ALIGN.CENTER
            run = p.runs[0]
            run.font.size = Pt(44)
            run.font.bold = True
            run.font.color.rgb = title_color

            if subtitle:
                stxb = slide.shapes.add_textbox(Inches(1), Inches(4.2), Inches(11), Inches(0.8))
                stf = stxb.text_frame
                stf.text = subtitle
                sp = stf.paragraphs[0]
                sp.alignment = PP_ALIGN.CENTER
                srun = sp.runs[0]
                srun.font.size = Pt(22)
                srun.font.color.rgb = text_color

            # Content slides
            for slide_data in slides:
                sl = prs.slides.add_slide(prs.slide_layouts[6])
                set_bg(sl)

                # Slide title
                ttxb = sl.shapes.add_textbox(Inches(0.7), Inches(0.4), Inches(11.9), Inches(1))
                ttf = ttxb.text_frame
                ttf.text = slide_data["title"] if isinstance(slide_data, dict) else slide_data.title
                tp = ttf.paragraphs[0]
                trun = tp.runs[0]
                trun.font.size = Pt(32)
                trun.font.bold = True
                trun.font.color.rgb = title_color

                # Accent line
                from pptx.util import Pt as Pt2
                line = sl.shapes.add_connector(1, Inches(0.7), Inches(1.55), Inches(12.3), Inches(1.55))
                line.line.color.rgb = accent_color
                line.line.width = Pt2(2)

                # Content
                content_text = slide_data["content"] if isinstance(slide_data, dict) else slide_data.content
                ctxb = sl.shapes.add_textbox(Inches(0.7), Inches(1.7), Inches(11.9), Inches(5.2))
                ctf = ctxb.text_frame
                ctf.word_wrap = True

                lines = content_text.split("\n")
                for i, line_text in enumerate(lines):
                    if i == 0:
                        p = ctf.paragraphs[0]
                    else:
                        p = ctf.add_paragraph()
                    p.text = line_text.strip("•-– ") if line_text.strip().startswith(("•", "-", "–")) else line_text
                    is_bullet = line_text.strip().startswith(("•", "-", "–"))
                    run = p.runs[0] if p.runs else p.add_run()
                    run.font.size = Pt(18)
                    run.font.color.rgb = text_color
                    if is_bullet:
                        p.level = 1

                # Speaker notes
                notes_text = slide_data.get("notes", "") if isinstance(slide_data, dict) else (slide_data.notes or "")
                if notes_text:
                    notes_slide = sl.notes_slide
                    notes_slide.notes_text_frame.text = notes_text

            prs.save(str(filepath))

            # Try to open on macOS
            try:
                import subprocess, sys
                if sys.platform == "darwin":
                    subprocess.Popen(["open", str(filepath)])
            except Exception:
                pass

            return f"Presentation created: {filepath}\n{len(slides)} content slides + cover slide."
        except ImportError:
            return "python-pptx not installed. Run: pip install python-pptx"
        except Exception as e:
            return f"Presentation creation failed: {str(e)}"


class CreateMarkdownReportInput(BaseModel):
    filename: str = Field(description="Filename for the report (e.g., 'trend-report.md')")
    title: str = Field(description="Report title")
    sections: list[dict] = Field(description="List of sections with 'heading' and 'content' keys")


class CreateMarkdownReportTool(BaseTool):
    name: str = "create_report"
    description: str = "Create a formatted Markdown report and save it to the workspace. Great for research reports, briefs, and summaries."
    args_schema: type[BaseModel] = CreateMarkdownReportInput

    def _run(self, filename: str, title: str, sections: list) -> str:
        try:
            ws = _ensure_workspace()
            if not filename.endswith(".md"):
                filename += ".md"
            filepath = ws / filename

            from datetime import datetime
            lines = [f"# {title}", f"*Generated by The Boss — {datetime.now().strftime('%Y-%m-%d %H:%M')}*", ""]
            for section in sections:
                heading = section.get("heading", "Section")
                content = section.get("content", "")
                lines.append(f"## {heading}")
                lines.append(content)
                lines.append("")

            filepath.write_text("\n".join(lines), encoding="utf-8")
            return f"Report saved: {filepath}"
        except Exception as e:
            return f"Report creation failed: {str(e)}"
